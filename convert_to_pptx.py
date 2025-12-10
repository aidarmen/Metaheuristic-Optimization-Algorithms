"""
Convert Slidev PNG exports to PowerPoint presentation
"""

import os
from pptx import Presentation
from pptx.util import Inches
import glob

def create_pptx_from_images(image_dir='presentation/slides-export', output_file='presentation/slides.pptx'):
    """
    Create a PowerPoint presentation from PNG images.
    
    Parameters:
    -----------
    image_dir : str
        Directory containing PNG slide images
    output_file : str
        Output PPTX filename
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get all PNG files sorted by name
    image_files = sorted(glob.glob(os.path.join(image_dir, '*.png')))
    
    if not image_files:
        print(f"No PNG files found in {image_dir}")
        return
    
    print(f"Found {len(image_files)} slide images")
    
    for i, img_path in enumerate(image_files, 1):
        print(f"Adding slide {i}: {os.path.basename(img_path)}")
        
        # Create blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add image to fill entire slide
        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_path, left, top, 
                                      width=prs.slide_width, 
                                      height=prs.slide_height)
    
    # Save presentation
    prs.save(output_file)
    print(f"\nPowerPoint presentation saved as: {output_file}")
    print(f"  Total slides: {len(image_files)}")

if __name__ == "__main__":
    create_pptx_from_images()

